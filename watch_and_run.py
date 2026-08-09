"""MinitAI engine: transcription, summarisation and document generation.

Imported by web.py (the local app) and by the hosted version, which reuses the
document generators without the local AI stack.
"""
import os
import re
import time
import json
import imageio_ffmpeg
os.environ["PATH"] = os.path.dirname(imageio_ffmpeg.get_ffmpeg_exe()) + os.pathsep + os.environ["PATH"]

import requests
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

import logging
import sqlite3
import sys

BASE_DIR = os.path.dirname(sys.executable) if getattr(sys, "frozen", False) else os.path.dirname(os.path.abspath(__file__))
DATA_DIR = os.path.join(os.environ.get("APPDATA", BASE_DIR), "MinitAI")
try:
    os.makedirs(DATA_DIR, exist_ok=True)
    _t = os.path.join(DATA_DIR, ".wtest")
    open(_t, "w").close(); os.remove(_t)
except OSError:
    import tempfile
    DATA_DIR = os.path.join(tempfile.gettempdir(), "MinitAI")
    os.makedirs(DATA_DIR, exist_ok=True)
AUDIO_DIR = os.path.join(DATA_DIR, "audio")
OUTPUT_DIR = os.path.join(DATA_DIR, "output")


def _desktop_folder():
    """Best-effort path to the user's Desktop (handles OneDrive-redirected Desktops)."""
    home = os.path.expanduser("~")
    candidates = [
        os.path.join(home, "Desktop"),
        os.path.join(home, "OneDrive", "Desktop"),
        os.path.join(os.environ.get("USERPROFILE", home), "Desktop"),
    ]
    for c in candidates:
        if os.path.isdir(c):
            return c
    return None


def _safe_folder_name(name, fallback="Meeting"):
    """A meeting title turned into something Windows will accept as a folder."""
    bad = '<>:"/\\|?*'
    out = "".join(("-" if ch in bad else ch) for ch in (name or "")).strip(" .")
    out = " ".join(out.split())          # collapse whitespace
    return (out[:70].strip(" .-") or fallback)


def save_to_desktop(*paths, folder=None):
    """Copy generated files to the Desktop.

    With `folder`, everything for one meeting lands in its own dated subfolder
    so the Desktop copy is browsable months later instead of a flat pile of
    timestamped filenames.
    """
    if not _cfg.get("save_to_desktop", True):
        return None
    try:
        import shutil
        desk = _desktop_folder()
        if not desk:
            return None
        target = os.path.join(desk, "MinitAI Documents")
        if folder:
            target = os.path.join(target, _safe_folder_name(folder))
        os.makedirs(target, exist_ok=True)
        # Inside a per-meeting folder the timestamp in each filename is just
        # noise - the folder already carries the date. Give them plain names.
        _nice = {".docx": "Minutes.docx", ".pptx": "Slides.pptx",
                 ".pdf": "Minutes.pdf", ".txt": "Transcript.txt"}
        saved = []
        for p in paths:
            if p and os.path.exists(p):
                base = os.path.basename(p)
                if folder:
                    ext = os.path.splitext(base)[1].lower()
                    base = _nice.get(ext, base)
                dest = os.path.join(target, base)
                # avoid clobber: add (n) if exists
                if os.path.exists(dest):
                    b, e = os.path.splitext(dest)
                    n = 1
                    while os.path.exists(f"{b} ({n}){e}"):
                        n += 1
                    dest = f"{b} ({n}){e}"
                shutil.copy2(p, dest)
                saved.append(dest)
        return target if saved else None
    except Exception:
        return None


LOG_PATH = os.path.join(DATA_DIR, "error_log.txt")
if os.path.exists(LOG_PATH) and os.path.getsize(LOG_PATH) > 5 * 1024 * 1024:
    os.replace(LOG_PATH, LOG_PATH + ".old")

APP_VERSION = "1.8.0"

def cleanup_old_files(days=30):
    cutoff = time.time() - days * 86400
    for d in (AUDIO_DIR,):
        if not os.path.isdir(d):
            continue
        for f in os.listdir(d):
            p = os.path.join(d, f)
            try:
                if os.path.isfile(p) and os.path.getmtime(p) < cutoff:
                    os.remove(p)
            except OSError:
                pass

def check_disk_space(min_gb=2):
    import shutil as _sh
    free_gb = _sh.disk_usage(DATA_DIR).free / (1024**3)
    if free_gb < min_gb:
        raise RuntimeError(f"Low disk space ({free_gb:.1f}GB free). Free up space and retry.")
logging.basicConfig(filename=LOG_PATH, level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(module)s:%(lineno)d %(message)s")
logging.getLogger().addHandler(logging.StreamHandler())
DB_PATH = os.path.join(DATA_DIR, "meetings.db")

def log_meeting(title, date, docx_path, pptx_path):
    try:
        conn = sqlite3.connect(DB_PATH)
        conn.execute("""CREATE TABLE IF NOT EXISTS meetings
            (id INTEGER PRIMARY KEY AUTOINCREMENT, title TEXT, date TEXT,
             docx_path TEXT, pptx_path TEXT, created_at TEXT DEFAULT (datetime('now')))""")
        conn.execute("INSERT INTO meetings (title,date,docx_path,pptx_path) VALUES (?,?,?,?)",
                     (title, date, docx_path, pptx_path))
        conn.commit()
        conn.close()
    except sqlite3.Error as e:
        logging.error(f"DB log failed: {e}")
AUDIO_EXTS = (".wav", ".mp3", ".m4a", ".mp4", ".aac", ".flac", ".ogg", ".wma", ".webm",
              ".mkv", ".mov", ".opus", ".amr", ".3gp", ".3gpp", ".aiff", ".caf")

CONFIG_FILE = os.path.join(DATA_DIR, "config.json")
if not os.path.exists(CONFIG_FILE):
    bundled = os.path.join(BASE_DIR, "config.json")
    if os.path.exists(bundled):
        import shutil
        shutil.copy(bundled, CONFIG_FILE)
_cfg = {}
if os.path.exists(CONFIG_FILE):
    try:
        with open(CONFIG_FILE) as f:
            _cfg = json.load(f)
    except Exception:
        _cfg = {}

def _detect_hardware_tier():
    """Pick model tier from RAM, CPU cores, and free memory. Returns (whisper, ollama, tier, info)."""
    ram_gb = 8.0
    avail_gb = 8.0
    cores = 4
    try:
        import psutil
        vm = psutil.virtual_memory()
        ram_gb = vm.total / (1024**3)
        avail_gb = vm.available / (1024**3)
        cores = psutil.cpu_count(logical=False) or psutil.cpu_count() or 4
    except Exception:
        try:
            cores = os.cpu_count() or 4
        except Exception:
            pass

    info = {"ram_gb": round(ram_gb, 1), "avail_gb": round(avail_gb, 1), "cores": cores}

    # Use free RAM as the real constraint (a 16GB machine with 3GB free will thrash).
    # Cores gate the top tier only; low cores step down one level, not to the floor.
    if ram_gb >= 16 and avail_gb >= 8 and cores >= 4:
        return ("small", "qwen2.5:7b", "high", info)
    # Whisper and the LLM are compute-bound on CPU, so cores gate every tier,
    # not just the top one. A 2-core machine given qwen2.5:3b crawls and looks
    # like a crash; that was the "friend's laptop" failure.
    if ram_gb >= 8 and avail_gb >= 4 and cores >= 4:
        return ("base", "qwen2.5:3b", "balanced", info)
    if ram_gb >= 16 and avail_gb >= 8:
        # Plenty of RAM but few cores: keep the better Whisper (RAM is free),
        # drop to the small LLM, and use the light resource policy so the
        # machine stays responsive.
        return ("base", "qwen2.5:1.5b", "light", info)
    return ("tiny", "qwen2.5:1.5b", "light", info)

_auto_whisper, _auto_ollama, _tier, _hw_info = _detect_hardware_tier()

# Fast mode: use a smaller/faster analysis model even on strong PCs.
# 7b is noticeably slow on CPU; 3b is much faster with good-enough minutes.
if _cfg.get("fast_mode", False) and _auto_ollama == "qwen2.5:7b":
    _auto_ollama = "qwen2.5:3b"

# user config overrides auto-detection; otherwise use hardware-appropriate defaults
MODEL_SIZE = _cfg.get("whisper_model_size") or _auto_whisper

# A converted CTranslate2 folder (e.g. a Malaysian-finetuned Whisper) may be used
# instead of a stock model name. The folder must be complete and local — see
# README "Malaysian model". Silently ignored if the path does not exist, so a
# bad config can never stop the app from starting.
def _find_local_whisper():
    """A converted CTranslate2 model dropped into %APPDATA%\\MinitAI\\models\\whisper
    is used automatically - no config editing. Installed by
    "More Options\\Use Malaysian Model.bat"."""
    auto = os.path.join(DATA_DIR, "models", "whisper")
    if os.path.exists(os.path.join(auto, "model.bin")):
        return auto
    return ""


WHISPER_MODEL_PATH = (_cfg.get("whisper_model_path") or "").strip() or _find_local_whisper()
if WHISPER_MODEL_PATH:
    if os.path.isdir(WHISPER_MODEL_PATH) and \
       os.path.exists(os.path.join(WHISPER_MODEL_PATH, "model.bin")):
        MODEL_SIZE = WHISPER_MODEL_PATH
        logging.info(f"using local Whisper model folder: {WHISPER_MODEL_PATH}")
    else:
        logging.warning(f"whisper_model_path not usable, ignoring: {WHISPER_MODEL_PATH}")
NAVY = RGBColor(0x1E, 0x27, 0x61)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x21, 0x21, 0x21)

OUTPUT_LANG = _cfg.get("output_language", "auto")  # auto|en|ms|zh
MEETING_TYPE = _cfg.get("meeting_type", "general")  # general|staff|committee|academic

_LANG_RULE = {
    "en": "Write all output in English.",
    "ms": "Tulis semua output dalam Bahasa Melayu.",
    "zh": "所有输出使用中文。",
}.get(OUTPUT_LANG, "Match the transcript's dominant language.")

_TYPE_RULE = {
    "staff": "This is a staff meeting: emphasize operational decisions and task assignments.",
    "committee": "This is a committee meeting: emphasize motions, votes, and formal resolutions.",
    "academic": "This is an academic meeting: emphasize curriculum, research, and student matters.",
}.get(MEETING_TYPE, "")

SYSTEM_PROMPT = _LANG_RULE + " " + _TYPE_RULE + """
You write official meeting minutes. The user message is a raw, automatic
transcript of a real meeting - it contains mis-heard words, false starts and
mixed Malay/English. Your job is to turn it into accurate minutes.

THE ONE RULE: write only what was actually said.
- Never invent a name, date, number, department or decision.
- If the transcript does not say it, leave that field empty. An empty field is
  correct. A guessed field is a serious error - these minutes are official.
- If a name is garbled, leave it out rather than guessing the spelling.
- Do not add polite filler, recommendations or conclusions of your own.

WHAT GOES WHERE
- meeting_title: the meeting's real name if stated (e.g. "Mesyuarat JK
  Pascasiswazah Bil 1/2026"), otherwise a short factual description.
- date / time / location: only if spoken aloud. Otherwise "".
- attendees: only people clearly present or greeted. Not people merely mentioned.
- agenda_items: one per topic actually discussed.
    topic      = what it was about, a few words
    discussion = what was said, 1-3 sentences, neutral and factual
    decision   = what was agreed. If nothing was agreed, write "No decision recorded".
- action_items: anything that still has to be DONE after this meeting.
    A decision almost always creates one. "Meluluskan X" means someone now has
    to carry X out; "borang perlu dihantar" is an action even though no name was
    said. Record it. An action with nobody named is still an action - the
    committee needs the list, and the owner can be filled in by hand.
    Do NOT leave action_items empty just because the meeting never said
    "you do this". Read the decisions and the "perlu / need to / should"
    sentences and write down what each one obliges someone to do.
    task = the thing to be done, phrased as an instruction.
    owner = the person named, or empty if nobody was named.
    deadline = only if a date or timeframe was stated.
- key_points: the substantive facts raised.
- key_takeaways: what matters going forward.
- important_notes: warnings, constraints, follow-ups, anything flagged.
- activities: events, visits, programmes mentioned.
- slides: a short deck. title_slide plus one content_slide per major topic.
  Bullets under 15 words, at most 5 per slide. Do not simply copy agenda_items -
  group related material so the deck reads well on its own.
- theme: pick primary_hex/accent_hex (6 hex chars, no #) to suit the subject -
  finance navy/gold, academic navy/gold, tech teal/cyan, health blue/green.
  mood = one word.

Write in the same language the meeting was conducted in unless told otherwise.
Keep the register formal, as official minutes are written."""

import os as _os
# Physical cores. os.cpu_count() returns LOGICAL cores, so on any hyperthreaded
# machine _threads was set above the number of real cores — CTranslate2 gets
# slower, not faster, when oversubscribed, and the OS is starved.
# _detect_hardware_tier already counts physical cores; reuse that number.
_cpu = _hw_info.get("cores") or _os.cpu_count() or 4

# --- Resource policy by hardware tier -----------------------------------
# On low-RAM machines (4GB) keeping models resident causes swap thrashing:
# the laptop is fine at first, then degrades badly. So on "light" we free
# memory aggressively and leave CPU headroom for the OS.
if _tier == "light":
    _threads = max(1, _cpu // 2)          # leave cores for the OS -> UI stays responsive
    _KEEP_ALIVE = "0"                     # unload the AI model immediately after use
    _FREE_WHISPER_BETWEEN = True          # release Whisper while the LLM runs
elif _tier == "balanced":
    _threads = max(1, _cpu - 1)
    _KEEP_ALIVE = "2m"
    _FREE_WHISPER_BETWEEN = False
else:                                      # high
    _threads = max(1, _cpu - 1)
    _KEEP_ALIVE = "5m"
    _FREE_WHISPER_BETWEEN = False

logging.info(f"tier={_tier} threads={_threads} keep_alive={_KEEP_ALIVE}")

# Whisper is loaded lazily so it isn't holding RAM before it's needed,
# and can be released again on low-RAM machines.
_whisper = None

def get_whisper():
    global _whisper
    if _whisper is None:
        from faster_whisper import WhisperModel     # imported on first use only
        logging.info(f"loading Whisper ({MODEL_SIZE}) with {_threads} threads")
        _whisper = WhisperModel(MODEL_SIZE, device="cpu", compute_type="int8", cpu_threads=_threads)
    return _whisper

def free_whisper():
    """Release the Whisper model and reclaim its RAM (low-memory machines)."""
    global _whisper
    if _whisper is not None:
        _whisper = None
        import gc
        gc.collect()
        logging.info("Whisper released to free RAM")

def unload_ollama_model():
    """Ask Ollama to drop the model from RAM right away."""
    try:
        requests.post(OLLAMA_URL + "/api/generate",
                      json={"model": (_resolved_model or OLLAMA_MODEL), "prompt": "", "keep_alive": 0},
                      timeout=15)
        logging.info("Ollama model unloaded")
    except Exception:
        pass
OLLAMA_MODEL = _cfg.get("ollama_model") or _auto_ollama
OLLAMA_URL = _cfg.get("ollama_url", "http://localhost:11434")


# --- Long-audio handling -------------------------------------------------
# Uploaded files longer than LONG_AUDIO_SECONDS are split into SEGMENT_SECONDS
# chunks before transcription, so a 55-minute upload can't hang the machine.
# (Live recording already segments in the browser; this covers UPLOADS.)
SEGMENT_SECONDS = 15 * 60      # 15 minutes per chunk — matches live recording
LONG_AUDIO_SECONDS = 16 * 60   # anything over ~16 min gets chunked


def _ffmpeg_exe():
    return imageio_ffmpeg.get_ffmpeg_exe()


def get_audio_duration(path):
    """Duration in seconds, or None if it can't be determined."""
    import subprocess as _sp, re as _re
    try:
        p = _sp.run([_ffmpeg_exe(), "-i", path],
                    capture_output=True, text=True, timeout=60)
        m = _re.search(r"Duration:\s*(\d+):(\d+):(\d+\.?\d*)", p.stderr or "")
        if m:
            h, mi, s = int(m.group(1)), int(m.group(2)), float(m.group(3))
            return h * 3600 + mi * 60 + s
    except Exception as e:
        logging.warning(f"duration probe failed: {e}")
    # A file recorded in the browser has no duration in its header - MediaRecorder
    # writes the stream as it goes and never goes back to fill it in. Decoding is
    # the only way to find out, and it is fast: a two-hour file takes a few
    # seconds. Without this a browser recording of any length is charged one
    # minute and skips the too-long check.
    try:
        p = _sp.run([_ffmpeg_exe(), "-i", path, "-f", "null", "-"],
                    capture_output=True, text=True, timeout=300)
        hits = _re.findall(r"time=\s*(\d+):(\d+):(\d+\.?\d*)", p.stderr or "")
        if hits:
            h, mi, s = hits[-1]
            secs = int(h) * 3600 + int(mi) * 60 + float(s)
            if secs > 0:
                return secs
    except Exception as e:
        logging.warning(f"duration decode failed: {e}")
    return None


def split_audio(path, out_dir, seconds=SEGMENT_SECONDS):
    """Split an audio/video file into ~`seconds`-long WAV chunks (16kHz mono, what
    Whisper wants). Returns an ordered list of chunk paths.
    Used so a long UPLOADED file doesn't overwhelm the machine in one pass."""
    import subprocess as _sp, glob as _glob
    os.makedirs(out_dir, exist_ok=True)
    pattern = os.path.join(out_dir, "chunk_%04d.wav")
    cmd = [
        _ffmpeg_exe(), "-hide_banner", "-loglevel", "error", "-y",
        "-i", path,
        "-vn",                     # drop video if present
        "-ac", "1", "-ar", "16000",  # mono 16kHz = Whisper's native format
        "-f", "segment", "-segment_time", str(int(seconds)),
        "-reset_timestamps", "1",
        pattern,
    ]
    p = _sp.run(cmd, capture_output=True, text=True, timeout=1800)
    chunks = sorted(_glob.glob(os.path.join(out_dir, "chunk_*.wav")))
    if not chunks:
        raise RuntimeError(f"Could not split the audio file. {(p.stderr or '')[:200]}")
    return chunks


# --- Transcription accuracy settings ------------------------------------
# Greedy decoding (beam_size=1) is the fastest option but measurably worse on
# mixed Malay/English speech, where the decoder has to choose between languages
# mid-sentence. Machines that can afford a beam should use one.
_TIER_BEAM = {"high": 5, "balanced": 3, "light": 1}
try:
    BEAM_SIZE = max(1, int(_cfg.get("beam_size") or _TIER_BEAM.get(_tier, 1)))
except (TypeError, ValueError):
    BEAM_SIZE = _TIER_BEAM.get(_tier, 1)

# Pin the language instead of auto-detecting. Auto-detect ran independently on
# every 15-minute chunk, so one English-heavy stretch could flip a Malay meeting
# mid-way and corrupt that chunk. "" / null = detect once, then lock.
TRANSCRIBE_LANG = (_cfg.get("transcribe_language") or "").strip() or None

# Seeds the decoder with names, acronyms and course codes it would otherwise
# mishear. Whisper only reads the last ~200 tokens, so keep it short.
INITIAL_PROMPT = (_cfg.get("initial_prompt") or "").strip() or None

_detected_lang = None   # locked on the first chunk of each job

logging.info(f"transcribe: beam={BEAM_SIZE} lang={TRANSCRIBE_LANG or 'auto-lock'} "
             f"prompt={'yes' if INITIAL_PROMPT else 'no'}")


# --- Engine routing: local CPU vs Groq cloud ------------------------------
# Machines differ enormously (a GTX-class desktop vs a 2-core laptop), so we do
# not guess from specs. We MEASURE how fast this machine actually transcribes
# and remember it, then route each job on the estimate.
try:
    import cloud as _cloud
except Exception:                                   # cloud.py missing = local only
    _cloud = None

ENGINE_MODE = (_cfg.get("engine") or "auto").lower()      # auto | local | cloud
# If the local engine would take longer than this, prefer the cloud.
CLOUD_IF_SLOWER_THAN = float(_cfg.get("cloud_if_slower_than_seconds", 600))
# Meetings the user marks confidential never leave the machine.
FORCE_LOCAL_ALWAYS = bool(_cfg.get("always_private", False))

_PERF_PATH = os.path.join(DATA_DIR, "perf.json")
# Starting guess until this machine has measured itself once.
_DEFAULT_RTF = {"light": 3.0, "balanced": 1.5, "high": 0.8}


def _load_rtf():
    """Seconds of CPU time per second of audio, measured on THIS machine."""
    try:
        with open(_PERF_PATH, encoding="utf-8") as f:
            d = json.load(f)
        r = float(d.get("rtf", 0))
        if 0.01 < r < 200:
            return r
    except Exception:
        pass
    return _DEFAULT_RTF.get(_tier, 2.0)


def record_rtf(audio_seconds, elapsed_seconds):
    """Remember how fast this machine really is (rolling average)."""
    if not audio_seconds or audio_seconds <= 0 or elapsed_seconds <= 0:
        return
    new = elapsed_seconds / audio_seconds
    try:
        old, n = _load_rtf(), 0
        try:
            with open(_PERF_PATH, encoding="utf-8") as f:
                n = int(json.load(f).get("samples", 0))
        except Exception:
            n = 0
        blended = new if n == 0 else (old * min(n, 5) + new) / (min(n, 5) + 1)
        with open(_PERF_PATH, "w", encoding="utf-8") as f:
            json.dump({"rtf": round(blended, 3), "samples": n + 1,
                       "tier": _tier, "threads": _threads}, f)
        logging.info(f"perf: this machine runs at {blended:.2f}x "
                     f"(1 hour of audio ~ {blended * 60:.0f} min)")
    except Exception as e:
        logging.info(f"perf: could not save timing: {type(e).__name__}")


def estimate_local_seconds(duration):
    return (duration or 0) * _load_rtf()


def choose_engine(duration=None, private=False):
    """Return 'cloud' or 'local' with a reason, for this specific job."""
    if FORCE_LOCAL_ALWAYS or private:
        return "local", "this meeting is marked private"
    if ENGINE_MODE == "local":
        return "local", "set to local in settings"
    have_cloud = bool(_cloud and _cloud.available(DATA_DIR))
    if ENGINE_MODE == "cloud":
        if have_cloud:
            return "cloud", "set to online in settings"
        return "local", "online mode is unavailable right now"
    # auto
    if not have_cloud:
        return "local", "no internet or no key"
    est = estimate_local_seconds(duration)
    if duration and est > CLOUD_IF_SLOWER_THAN:
        return "cloud", (f"this PC would need about {est / 60:.0f} min; "
                         f"online is much faster")
    if not duration:
        return "local", "length unknown"
    return "local", f"this PC can do it in about {max(1, est / 60):.0f} min"


def transcribe(audio_path, progress=None):
    """Transcribe an audio file. Long files are split into chunks first so a big
    upload (e.g. a 55-minute meeting) can't hang the machine."""
    global _detected_lang
    _detected_lang = None          # never inherit a language from the last meeting
    check_disk_space()
    if len(audio_path) > 250:
        raise RuntimeError("File path too long. Move file to a shorter path.")
    if not os.path.exists(audio_path) or os.path.getsize(audio_path) < 1024:
        raise RuntimeError("Audio file missing or too small/corrupt.")

    duration = get_audio_duration(audio_path)

    engine, why = choose_engine(duration, private=_cfg.get("_private_job", False))
    logging.info(f"engine: {engine} ({why})")
    if engine == "cloud":
        try:
            t0 = time.time()
            text = _cloud.transcribe(audio_path, DATA_DIR, _ffmpeg_exe(),
                                     language=TRANSCRIBE_LANG,
                                     prompt=INITIAL_PROMPT,
                                     duration=duration, progress=progress)
            logging.info(f"cloud transcribe took {time.time() - t0:.0f}s")
            return text
        except Exception as e:
            # Never let a cloud problem stop the meeting from being processed.
            logging.warning(f"cloud failed ({type(e).__name__}), using this PC instead")

    _t_start = time.time()
    if duration and duration > LONG_AUDIO_SECONDS:
        logging.info(f"transcribe: {duration/60:.1f} min audio — splitting into chunks")
        work = os.path.join(AUDIO_DIR, f"_chunks_{int(time.time())}")
        try:
            chunks = split_audio(audio_path, work)
            logging.info(f"transcribe: {len(chunks)} chunk(s)")
            parts, failed = [], []
            for i, ch in enumerate(chunks, 1):
                if progress:
                    try:
                        progress(i, len(chunks))
                    except Exception:
                        pass
                logging.info(f"transcribe: chunk {i}/{len(chunks)}")
                try:
                    parts.append(_transcribe_one(ch))
                except Exception as e:
                    # A swallowed chunk is a silent 15-minute hole: the minutes
                    # still look complete. Make the gap visible instead.
                    failed.append(i)
                    logging.warning(f"chunk {i} failed: {e}")
                    parts.append(
                        f"\n[MinitAI: part {i} of {len(chunks)} could not be transcribed "
                        f"— roughly {SEGMENT_SECONDS // 60} minutes are missing here]\n")
                finally:
                    try:
                        os.remove(ch)      # free disk as we go
                    except Exception:
                        pass
                    import gc
                    gc.collect()           # reclaim RAM between chunks (matters on 4GB)
            if failed and len(failed) * 2 >= len(chunks):
                raise RuntimeError(
                    f"Transcription failed on {len(failed)} of {len(chunks)} parts, "
                    "so most of the meeting would be missing from the minutes. "
                    "Check free disk space and available memory, then try again.")
            if failed:
                logging.warning(f"transcribe: {len(failed)}/{len(chunks)} chunks failed: {failed}")
            text = " ".join(t for t in parts if t.strip())
            if not text.strip():
                raise RuntimeError("No speech detected in audio.")
            record_rtf(duration, time.time() - _t_start)
            return text
        finally:
            import shutil
            shutil.rmtree(work, ignore_errors=True)

    _out = _transcribe_one(audio_path)
    record_rtf(duration, time.time() - _t_start)
    return _out


def _transcribe_one(audio_path):
    """Transcribe a single (already short enough) audio file.
    vad_filter skips silence, saving time on real recordings."""
    global _detected_lang
    lang = TRANSCRIBE_LANG or _detected_lang
    segments, info = get_whisper().transcribe(
        audio_path, language=lang, vad_filter=True,
        beam_size=BEAM_SIZE, best_of=BEAM_SIZE,
        condition_on_previous_text=False,
        initial_prompt=INITIAL_PROMPT,
    )
    text = " ".join(seg.text.strip() for seg in segments)
    if lang is None:
        # Lock whatever the first chunk detected so later chunks can't flip
        # language mid-meeting.
        detected = getattr(info, "language", None)
        if detected:
            _detected_lang = detected
            logging.info(f"language locked to '{detected}' for the rest of this file")
    if not text.strip():
        raise RuntimeError("No speech detected in audio.")
    return text


import subprocess

def ensure_ollama_ready():
    ollama_exe = os.path.join(os.environ.get("LOCALAPPDATA", ""), "Programs", "Ollama", "ollama.exe")
    if not os.path.exists(ollama_exe):
        ollama_exe = "ollama"

    try:
        requests.get(OLLAMA_URL, timeout=3)
    except requests.exceptions.ConnectionError:
        try:
            subprocess.Popen([ollama_exe, "serve"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        except FileNotFoundError:
            raise RuntimeError("Ollama not installed or not found. Reinstall MinitAI.")
        for _ in range(10):
            time.sleep(1)
            try:
                requests.get(OLLAMA_URL, timeout=3)
                break
            except requests.exceptions.ConnectionError:
                continue
        else:
            raise RuntimeError("Ollama failed to start. Open Ollama app manually.")

    tags = requests.get(OLLAMA_URL + "/api/tags", timeout=5).json()
    models = [m["name"] for m in tags.get("models", [])]
    if not any(OLLAMA_MODEL in m for m in models):
        try:
            import psutil
            ram_gb = psutil.virtual_memory().total / (1024**3)
            if ram_gb < 8:
                logging.warning(f"Low RAM ({ram_gb:.1f}GB) — model may fail to load.")
        except ImportError:
            pass
        try:
            subprocess.run([ollama_exe, "pull", OLLAMA_MODEL], check=False, timeout=600)
        except FileNotFoundError:
            pass
        except subprocess.TimeoutExpired:
            raise RuntimeError(f"Model download timed out. Check internet connection.")
        tags2 = requests.get(OLLAMA_URL + "/api/tags", timeout=5).json()
        models2 = [m["name"] for m in tags2.get("models", [])]
        if not any(OLLAMA_MODEL in m for m in models2):
            raise RuntimeError(f"Model {OLLAMA_MODEL} not available. Check internet, then retry.")


MAX_TRANSCRIPT_CHARS = 24000  # ~6k tokens, safe for 7B ctx
LONG_MEETING_CHARS = 14000    # beyond this (~15 min speech) we chunk instead of truncate
_in_chunk_pass = False        # guard so chunk passes don't recurse into analyze_long

_LANG_RULES = {
    "en": "Write all output in English.",
    "ms": "Tulis semua output dalam Bahasa Melayu.",
    "zh": "所有输出使用中文。",
    "auto": "Match the transcript's dominant language.",
}

# Words that carry no topic meaning, so two titles about the same thing still
# match when one says "Cadangan pemeriksa luar" and the other "Pelantikan
# pemeriksa luar dan dalam".
_TOPIC_STOP = {
    "dan", "atau", "bagi", "untuk", "pada", "dari", "daripada", "kepada",
    "yang", "dengan", "serta", "oleh", "itu", "ini", "adalah", "ialah",
    "tentang", "mengenai", "terhadap", "baru", "baharu",
    "the", "and", "for", "of", "to", "on", "in", "about", "new", "item",
}


# Malay builds words with affixes, so "tukar" and "pertukaran" are the same
# subject written two ways. Stripping them is what lets the two be recognised
# as one agenda item. peN-/meN- swallow the root's first letter, which is why
# "pemeriksa" has to become "periksa" and not "eriksa".
_MS_PREFIX = (("menge", ""), ("penge", ""), ("meng", "k"), ("peng", "k"),
              ("meny", "s"), ("peny", "s"), ("mem", "p"), ("pem", "p"),
              ("men", "t"), ("pen", "t"), ("ber", ""), ("ter", ""),
              ("per", ""), ("me", ""), ("pe", ""), ("di", ""), ("ke", ""))
_MS_SUFFIX = ("kannya", "annya", "nya", "kan", "an")


def _ms_stem(w):
    for pre, restore in _MS_PREFIX:
        if w.startswith(pre) and len(w) - len(pre) >= 3:
            rest = w[len(pre):]
            if restore and rest[0] in "aeiou":
                rest = restore + rest
            w = rest
            break
    for suf in _MS_SUFFIX:
        if w.endswith(suf) and len(w) - len(suf) >= 3:
            return w[:-len(suf)]
    return w


def _topic_key(text):
    """The words that carry the topic, reduced to their stems."""
    words = re.findall(r"[a-z0-9]+", (text or "").lower())
    return {_ms_stem(w) for w in words if len(w) > 2 and w not in _TOPIC_STOP}


def _fuzzy_overlap(ka, kb):
    """Shared stems, allowing for the transcript misspelling one of them.

    Whisper wrote the same word as "kolokium" once and "kolekium" the next
    time; exact matching treats those as two different subjects.
    """
    import difflib
    hits = 0
    left = set(kb)
    for a in ka:
        if a in left:
            hits += 1
            left.discard(a)
            continue
        near = next((b for b in left
                     if difflib.SequenceMatcher(None, a, b).ratio() >= 0.85), None)
        if near:
            hits += 1
            left.discard(near)
    return hits


def _same_topic(a, b):
    """Two agenda titles describing the same thing, not necessarily worded alike."""
    ka, kb = _topic_key(a), _topic_key(b)
    if not ka or not kb:
        return False
    inter = _fuzzy_overlap(ka, kb)
    if not inter:
        return False
    # one title's words wholly inside the other, or half the combined words shared
    return inter == min(len(ka), len(kb)) or inter / (len(ka) + len(kb) - inter) >= 0.5


def _real_decision(text):
    t = (text or "").strip().rstrip(".").lower()
    return bool(t) and t != "no decision recorded"


def _merge_agenda(items):
    """Fold chunk-level agenda items that cover the same topic into one.

    A two-hour meeting is analysed in chunks, and a topic raised across two
    chunks used to produce two near-identical entries. One real meeting came
    back with nineteen items covering about eleven topics.
    """
    out = []
    for it in items:
        if not isinstance(it, dict):
            continue
        topic = _as_text(it.get("topic")).strip()
        disc = _as_text(it.get("discussion")).strip()
        dec = _as_text(it.get("decision")).strip()
        hit = next((p for p in out if _same_topic(p["topic"], topic)), None)
        if hit is None:
            out.append({"topic": topic, "discussion": disc, "decision": dec})
            continue
        if len(disc) > len(hit["discussion"]):
            hit["discussion"] = disc
        if _real_decision(dec):
            if not _real_decision(hit["decision"]):
                hit["decision"] = dec
            elif dec.lower() not in hit["decision"].lower():
                hit["decision"] = hit["decision"].rstrip(".") + "; " + dec
        # The first title wins: it is where the topic was introduced, and it is
        # usually the fuller description. Later chunks tend to abbreviate.
    return out


def _merge_actions(items):
    """Same task recorded in two chunks is one task."""
    out, seen = [], set()
    for it in items:
        if not isinstance(it, dict):
            continue
        task = _as_text(it.get("task")).strip()
        if not task:
            continue
        key = " ".join(sorted(_topic_key(task)))
        if key in seen:
            continue
        seen.add(key)
        out.append({"task": task,
                    "owner": _as_text(it.get("owner")).strip(),
                    "deadline": _as_text(it.get("deadline")).strip()})
    return out


def _merge_analyses(parts):
    """Merge several chunk analyses into one meeting summary."""
    if not parts:
        return {}
    merged = {
        "meeting_title": "", "date": "", "time": "", "location": "",
        "attendees": [], "agenda_items": [], "activities": [],
        "key_points": [], "key_takeaways": [], "important_notes": [],
        "action_items": [],
        "theme": parts[0].get("theme", {"primary_hex": "1E2761", "accent_hex": "FFD500"}),
    }
    seen = {k: set() for k in ("attendees", "key_points", "key_takeaways", "important_notes", "activities")}
    for p in parts:
        for f in ("meeting_title", "date", "time", "location"):
            if not merged[f] and p.get(f):
                merged[f] = p[f]
        for f in ("attendees", "key_points", "key_takeaways", "important_notes", "activities"):
            for item in (p.get(f) or []):
                key = str(item).strip().lower()
                if key and key not in seen[f]:
                    seen[f].add(key)
                    merged[f].append(item)
        merged["agenda_items"].extend(p.get("agenda_items") or [])
        merged["action_items"].extend(p.get("action_items") or [])
    # Chunks overlap in subject matter, so fold duplicates before anything
    # downstream counts, numbers or renders them.
    merged["agenda_items"] = _merge_agenda(merged["agenda_items"])
    merged["action_items"] = _merge_actions(merged["action_items"])
    # rebuild slides from the merged content
    merged["slides"] = {
        "title_slide": {"title": merged["meeting_title"] or "Meeting Minutes",
                        "subtitle": merged.get("date", "")},
        "content_slides": [
            {"heading": a.get("topic", "Discussion"),
             "bullets": [b for b in [a.get("discussion"), a.get("decision")] if b]}
            for a in merged["agenda_items"][:12]
        ],
    }
    return merged


def analyze_long(transcript_text, lang=None, chunk_chars=None):
    """Analyze a long transcript by splitting it into chunks and merging the results.
    Prevents the blank-document failure on long meetings.
    Chunk size is tier-aware so each chunk fits the (RAM-capped) context window."""
    if chunk_chars is None:
        chunk_chars = {"light": 8000, "balanced": 12000}.get(_tier, 12000)
    # split on sentence boundaries near the chunk size
    chunks, cur = [], ""
    for sentence in transcript_text.replace("\n", " ").split(". "):
        if len(cur) + len(sentence) > chunk_chars and cur:
            chunks.append(cur.strip())
            cur = ""
        cur += sentence + ". "
    if cur.strip():
        chunks.append(cur.strip())

    logging.info(f"analyze_long: {len(transcript_text)} chars -> {len(chunks)} chunk(s)")
    global _in_chunk_pass
    parts = []
    _in_chunk_pass = True
    try:
        for i, ch in enumerate(chunks, 1):
            logging.info(f"analyze_long: chunk {i}/{len(chunks)}")
            try:
                parts.append(analyze(ch, lang=lang))
            except Exception as e:
                logging.warning(f"chunk {i} failed: {e}")
    finally:
        _in_chunk_pass = False
    if not parts:
        raise RuntimeError(
            "Couldn't summarize this meeting — the AI failed on every section. "
            'Try setting "fast_mode": true in config.json, or split the recording into shorter parts.'
        )
    return _merge_analyses(parts)


# --- Grammar-constrained output ------------------------------------------
# Ollama compiles a JSON Schema into a GBNF grammar, so the model physically
# CANNOT emit malformed or wrongly-shaped JSON. This removes the whole failure
# class behind bug #6 (blank Word/PPT) and the "AI returned invalid data" retry
# loop. Older Ollama builds ignore schemas, so analyze() falls back to
# format:"json" automatically if a schema request fails.
_STR = {"type": "string"}
ANALYSIS_SCHEMA = {
    "type": "object",
    "properties": {
        "meeting_title": _STR, "date": _STR, "time": _STR, "location": _STR,
        "attendees": {"type": "array", "items": _STR},
        "agenda_items": {"type": "array", "items": {
            "type": "object",
            "properties": {"topic": _STR, "discussion": _STR, "decision": _STR},
            "required": ["topic", "discussion", "decision"]}},
        "activities": {"type": "array", "items": _STR},
        "key_points": {"type": "array", "items": _STR},
        "key_takeaways": {"type": "array", "items": _STR},
        "important_notes": {"type": "array", "items": _STR},
        "action_items": {"type": "array", "items": {
            "type": "object",
            "properties": {"task": _STR, "owner": _STR, "deadline": _STR},
            "required": ["task", "owner", "deadline"]}},
        "theme": {"type": "object",
                  "properties": {"primary_hex": _STR, "accent_hex": _STR, "mood": _STR},
                  "required": ["primary_hex", "accent_hex", "mood"]},
        "slides": {"type": "object", "properties": {
            "title_slide": {"type": "object",
                            "properties": {"title": _STR, "subtitle": _STR},
                            "required": ["title", "subtitle"]},
            "content_slides": {"type": "array", "items": {
                "type": "object",
                "properties": {"heading": _STR, "bullets": {"type": "array", "items": _STR}},
                "required": ["heading", "bullets"]}}},
            "required": ["title_slide", "content_slides"]},
    },
    "required": ["meeting_title", "date", "time", "location", "attendees",
                 "agenda_items", "activities", "key_points", "key_takeaways",
                 "important_notes", "action_items", "theme", "slides"],
}

# Set false in config.json if your Ollama is old and schemas misbehave.
USE_SCHEMA = _cfg.get("strict_json", True)
_schema_supported = None       # None = untested, True/False after first attempt


# Titles are everywhere in a Malaysian meeting, so matching one proves nothing.
# "Datuk Seri Bayangan" slipped through the first version of this check purely
# because "datuk" appeared elsewhere in the transcript.
_HONORIFICS = {
    "dr", "drs", "prof", "profesor", "madya", "assoc", "associate",
    "datuk", "dato", "datin", "seri", "sri", "tan", "puan", "encik",
    "tuan", "cik", "haji", "hajah", "hj", "ir", "ts", "mr", "mrs", "ms",
    "bin", "binti", "bt", "al", "abd", "the", "and", "dan",
}


def _grounded(value, transcript_low):
    """True if a short extracted string plausibly came from the transcript.
    Guards against invented attendee names and owners on small models."""
    v = (value or "").strip()
    if len(v) < 3 or len(v) > 60:
        return True                      # too short/long to check safely
    low = v.lower()
    if low in transcript_low:
        return True
    words = [w.strip(".,'\"()") for w in low.replace(",", " ").split()]
    words = [w for w in words if len(w) > 2]
    if not words:
        return True
    # Only distinctive words count - a matching title is not evidence.
    distinctive = [w for w in words if w not in _HONORIFICS]
    if not distinctive:
        return True                      # nothing but titles; nothing to judge
    hits = sum(1 for w in distinctive if w in transcript_low)
    return hits >= max(1, (len(distinctive) + 1) // 2)


def _drop_hallucinations(data, transcript_text, keep=None):
    """Remove attendees and action owners that never appear in the transcript.

    `keep` is the roster the user typed in. Those names are exempt: someone who
    sat through a meeting without speaking is still an attendee, and deleting a
    name the user typed themselves is a worse error than keeping one the model
    invented. Without this, the guard silently undoes the roster.
    """
    low = (transcript_text or "").lower()
    if not low:
        return data
    trusted = {str(k).strip().lower() for k in (keep or []) if str(k).strip()}
    before = len(data.get("attendees") or [])
    data["attendees"] = [a for a in (data.get("attendees") or [])
                         if str(a).strip().lower() in trusted or _grounded(str(a), low)]
    dropped = before - len(data["attendees"])
    for it in (data.get("action_items") or []):
        if (isinstance(it, dict) and it.get("owner")
                and str(it["owner"]).strip().lower() not in trusted
                and not _grounded(str(it["owner"]), low)):
            it["owner"] = ""
            dropped += 1
    if dropped:
        logging.info(f"analysis: dropped {dropped} ungrounded name(s)")
    return data


# --- Model preference ----------------------------------------------------
# qwen2.5 shipped in 2024 and is the weakest link in summary quality. Newer
# small models are much better, and a Malaysian-tuned model is better still.
# We only ever SELECT AMONG MODELS ALREADY INSTALLED - never auto-pull - so a
# missing tag can't repeat bug #8 (installer grabbed the wrong/largest model).
# Grammar-constrained output (ANALYSIS_SCHEMA) guarantees valid JSON no matter
# which model wins, so the model only has to be good at content, not format.
# That is what makes a Malay-tuned model safe to prefer here.
_PREFERRED_MODELS = {
    "high":     ["Supa-AI/malaysian-llama-3.2-3b-instruct:q4_k_s",
                 "qwen3:8b", "qwen2.5:7b", "llama3.1:8b"],
    "balanced": ["Supa-AI/malaysian-llama-3.2-3b-instruct:q4_k_s",
                 "qwen3.5:4b", "qwen2.5:3b", "llama3.2:3b"],
    "light":    ["qwen2.5:1.5b", "llama3.2:1b", "gemma3:1b"],
}
_resolved_model = None


def installed_models():
    """Tags currently installed in Ollama. Empty list if it can't be reached."""
    try:
        r = requests.get(OLLAMA_URL + "/api/tags", timeout=5)
        return [m.get("name", "") for m in (r.json().get("models") or [])]
    except Exception:
        return []


def pick_model():
    """Best INSTALLED model for this tier, falling back to the configured one.
    Cached: Ollama is only asked once per run."""
    global _resolved_model
    if _resolved_model:
        return _resolved_model
    if _cfg.get("ollama_model"):
        _resolved_model = _cfg["ollama_model"]          # explicit user choice wins
        return _resolved_model
    have = installed_models()
    if have:
        base = {h.split(":")[0]: h for h in have}
        for want in _PREFERRED_MODELS.get(_tier, []):
            if want in have or want.split(":")[0] in base:
                _resolved_model = want if want in have else base[want.split(":")[0]]
                if _resolved_model != OLLAMA_MODEL:
                    logging.info(f"using better installed model: {_resolved_model}")
                return _resolved_model
    _resolved_model = OLLAMA_MODEL
    return _resolved_model


def _validate_analysis(data):
    """Reject an analysis that would produce a blank document."""
    if not isinstance(data, dict):
        return False
    title = (data.get("meeting_title") or "").strip()
    has_content = any([
        data.get("agenda_items"),
        data.get("key_points"),
        data.get("key_takeaways"),
        data.get("action_items"),
        data.get("activities"),
        data.get("important_notes"),
    ])
    return bool(title) or has_content


def analyze(transcript_text, retries=2, lang=None):
    global _schema_supported
    # Long meetings: process in chunks instead of truncating (which lost content
    # and caused blank documents). LONG_MEETING_CHARS ~ 15 min of speech.
    if len(transcript_text) > LONG_MEETING_CHARS and not _in_chunk_pass:
        return analyze_long(transcript_text, lang=lang)

    if len(transcript_text) > MAX_TRANSCRIPT_CHARS:
        logging.warning(f"Transcript {len(transcript_text)} chars — truncating to {MAX_TRANSCRIPT_CHARS}")
        transcript_text = transcript_text[:MAX_TRANSCRIPT_CHARS]

    # On low-RAM machines, release Whisper before loading the language model so
    # the two never sit in memory together (this is what causes 4GB laptops to
    # start swapping and lag badly part-way through a long meeting).
    if _FREE_WHISPER_BETWEEN:
        free_whisper()

    ensure_ollama_ready()

    system_prompt = SYSTEM_PROMPT
    if lang and lang in _LANG_RULES:
        system_prompt = _LANG_RULES[lang] + " " + _TYPE_RULE + SYSTEM_PROMPT.split("\n", 1)[1]

    # CRITICAL: the context window must fit the WHOLE transcript + system prompt + output,
    # or the model silently overflows and returns empty/garbage JSON (=> blank document).
    # ~3.5 chars per token is a safe estimate for mixed English/BM text.
    est_tokens = (len(transcript_text) + len(system_prompt)) // 3 + 2000  # +2000 for the JSON output
    # A big context costs a lot of RAM. Cap it per tier so a 4GB laptop can't thrash.
    ctx_cap = {"light": 8192, "balanced": 16384}.get(_tier, 32768)
    ctx = 4096
    for size in (4096, 8192, 16384, 32768):
        if est_tokens <= size:
            ctx = size
            break
    else:
        ctx = 32768
    ctx = min(ctx, ctx_cap)
    logging.info(f"analyze: {len(transcript_text)} chars -> num_ctx {ctx} (tier={_tier})")

    # Summary can also run in the cloud - far better output from a 70B model
    # than a 1.5B one on a weak laptop. Falls back to local on any failure.
    _eng, _why = choose_engine(None, private=_cfg.get("_private_job", False))
    if _eng == "cloud" or (ENGINE_MODE != "local" and _cloud
                           and _cloud.available(DATA_DIR)):
        try:
            data = _cloud.analyze(transcript_text, DATA_DIR, system_prompt,
                                  ANALYSIS_SCHEMA)
            if _validate_analysis(data):
                return _drop_hallucinations(data, transcript_text)
            logging.warning("cloud summary was empty, using this PC instead")
        except Exception as e:
            logging.warning(f"cloud summary failed ({type(e).__name__}), using this PC")

    ensure_ollama_ready()
    last_err = None
    for attempt in range(retries + 1):
        _payload = {
            "model": pick_model(),
            "system": system_prompt,
            "prompt": transcript_text,
            "stream": False,
            "format": (ANALYSIS_SCHEMA if (USE_SCHEMA and _schema_supported is not False)
                       else "json"),
            "keep_alive": _KEEP_ALIVE,
            "options": {
                "num_ctx": ctx,
                "num_predict": 2500,
            },
        }
        _sent_schema = _payload["format"] != "json"
        try:
            resp = requests.post(OLLAMA_URL + "/api/generate", timeout=900, json=_payload)
            if getattr(resp, "status_code", 200) >= 400 and _sent_schema:
                # Ollama older than v0.5 does not accept a JSON Schema. Drop to
                # plain JSON mode for the rest of this run instead of failing.
                _schema_supported = False
                logging.warning("Ollama rejected the JSON schema - using plain JSON mode")
                continue
            raw = resp.json().get("response", "")
        except Exception as e:
            last_err = f"Could not reach the AI engine: {e}"
            continue
        if _sent_schema:
            _schema_supported = True

        if not raw.strip():
            last_err = "The AI returned an empty response (the meeting may be too long for this model)."
            continue

        try:
            data = json.loads(raw)
        except json.JSONDecodeError:
            data = _extract_json_loose(raw)
            if data is None:
                # Describe the SHAPE of the bad response, never its content —
                # the response is derived from the meeting and must not reach logs.
                last_err = f"The AI returned invalid data ({_shape_of(raw)})."
                continue

        # Never let a blank document through.
        if not _validate_analysis(data):
            last_err = "The AI produced an empty summary — nothing usable was extracted."
            continue

        # Small models invent attendee names and owners. Drop anything that
        # does not actually appear in the transcript.
        data = _drop_hallucinations(data, transcript_text)

        # On low-RAM machines free the model right away so the laptop stays responsive.
        if _tier == "light" and not _in_chunk_pass:
            unload_ollama_model()
        return data

    raise RuntimeError(
        f"Couldn't summarize this meeting. {last_err or ''} "
        "Tip: very long meetings can overwhelm the local AI — try splitting the recording, "
        'or set "fast_mode": true in config.json.'
    )


def _shape_of(raw):
    """Describe a bad model response for debugging WITHOUT revealing its content.
    The response is derived from the user's meeting, so it must never be logged."""
    if raw is None:
        return "no response"
    s = str(raw)
    if not s.strip():
        return "empty response"
    title_key = '"meeting_title"'
    has_title = "has" if title_key in s else "no"
    first = s.lstrip()[:1]
    return (f"{len(s)} chars, starts with '{first}', "
            f"braces {s.count('{')}open/{s.count('}')}close, "
            f"{has_title} title key")


def _extract_json_loose(raw):
    """Salvage a JSON object from a messy model response."""
    import re as _re
    m = _re.search(r"\{.*\}", raw, _re.DOTALL)
    if not m:
        return None
    try:
        return json.loads(m.group(0))
    except json.JSONDecodeError:
        return None


from docx.shared import RGBColor as DocxRGBColor

def hex_to_docx_rgb(hexstr, fallback_hex="1E2761"):
    try:
        h = (hexstr or fallback_hex).lstrip('#')
        return DocxRGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
    except Exception:
        return DocxRGBColor(0x1E, 0x27, 0x61)


def hex_to_rgb(hexstr, fallback):
    try:
        h = hexstr.lstrip('#')
        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
    except Exception:
        return fallback


# --- Document labels ------------------------------------------------------
# Minutes written in Malay with English section headings ("ATTENDEES",
# "Decision:", "Task / Owner / Deadline") do not pass as an official document.
# Headings now follow the language of the minutes themselves.
_DOC_LABELS = {
    "en": {"attendees": "Attendees", "activities": "Activities",
           "matters": "Matters Discussed", "decision": "Decision: ",
           "key_points": "Key Points", "takeaways": "Key Takeaways",
           "notes": "Important Notes", "actions": "Action Items",
           "task": "Task", "owner": "Owner", "deadline": "Deadline",
           "date": "Date", "time": "Time", "location": "Location",
           "no_decision": "No decision recorded",
           "prepared_by": "Prepared by", "verified_by": "Verified and confirmed by",
           "name_line": "Name", "sig_date": "Date",
           "fallback_title": "Meeting Minutes"},
    "ms": {"attendees": "Kehadiran", "activities": "Aktiviti",
           "matters": "Perkara Dibincangkan", "decision": "Keputusan: ",
           "key_points": "Perkara Utama", "takeaways": "Rumusan",
           "notes": "Catatan Penting", "actions": "Tindakan",
           "task": "Tindakan", "owner": "Pegawai Bertanggungjawab",
           "deadline": "Tarikh Akhir",
           "date": "Tarikh", "time": "Masa", "location": "Tempat",
           "no_decision": "Tiada keputusan direkodkan",
           "prepared_by": "Disediakan oleh", "verified_by": "Disemak dan disahkan oleh",
           "name_line": "Nama", "sig_date": "Tarikh",
           "fallback_title": "Minit Mesyuarat"},
    "zh": {"attendees": "出席者", "activities": "活动",
           "matters": "讨论事项", "decision": "决定: ",
           "key_points": "重点", "takeaways": "总结",
           "notes": "重要事项", "actions": "行动事项",
           "task": "事项", "owner": "负责人", "deadline": "期限",
           "date": "日期", "time": "时间", "location": "地点",
           "no_decision": "未作出决定",
           "prepared_by": "记录人", "verified_by": "审核及确认人",
           "name_line": "姓名", "sig_date": "日期",
           "fallback_title": "会议记录"},
}

# Common Malay words used to detect the language of the minutes when the user
# left output_language on "auto".
_MS_HINTS = ("mesyuarat", "pelajar", "yang", "dan", "untuk", "kepada", "tidak",
             "adalah", "perlu", "telah", "dalam", "kertas", "tindakan",
             "keputusan", "bagi", "dengan", "ini", "itu")


def doc_labels(data=None):
    """Which language the document furniture should be written in."""
    lang = (OUTPUT_LANG or "auto").lower()
    if lang in _DOC_LABELS:
        return _DOC_LABELS[lang]
    # auto: judge from the generated content itself
    try:
        blob = json.dumps(data or {}, ensure_ascii=False).lower()
        if any(ord(ch) > 0x4DFF for ch in blob[:4000]):
            return _DOC_LABELS["zh"]
        hits = sum(1 for wrd in _MS_HINTS if f" {wrd} " in blob or f'"{wrd}' in blob)
        if hits >= 3:
            return _DOC_LABELS["ms"]
    except Exception:
        pass
    return _DOC_LABELS["en"]


def _as_dict(v, default=None):
    return v if isinstance(v, dict) else (default if default is not None else {})


def _as_list(v):
    if isinstance(v, list):
        return v
    if v in (None, "", {}):
        return []
    return [v]


def _as_text(v):
    """Flatten whatever a model produced into a plain string."""
    if isinstance(v, str):
        return v
    if v is None or isinstance(v, (dict, list)) and not v:
        return ""
    if isinstance(v, (int, float)):
        return str(v)
    if isinstance(v, list):
        return " ".join(_as_text(x) for x in v if x)
    if isinstance(v, dict):
        return " ".join(_as_text(x) for x in v.values() if x)
    return str(v)


_TITLES = {"prof", "profesor", "dr", "datuk", "dato", "datin", "tuan", "puan",
           "encik", "cik", "tn", "pn", "en", "madya", "ir", "ts", "haji", "hajah",
           "yb", "ybhg", "mr", "mrs", "ms", "miss"}


def _name_core(name):
    """The name with titles and punctuation stripped, for comparison only."""
    words = [w for w in re.findall(r"[A-Za-z']+", name)
             if w.lower().strip("'") not in _TITLES]
    return " ".join(words).lower()


def _clean_attendees(names):
    """One person, one spelling, and no bare titles.

    A transcript produced two attendees called "Prof." with no name attached,
    and wrote the same person as "Dr. Zuleti" and "Dr. Zulaiti".

    Only certain duplicates are merged. Names that merely LOOK alike are left
    alone and reported instead: quietly collapsing two colleagues into one
    would be a worse error in official minutes than listing a name twice, and
    only a human knows whether Hafizah and Hafizal are the same person.

    Returns (attendees, warnings).
    """
    import difflib
    out, warn = [], []
    for n in names:
        core = _name_core(n)
        if not core:
            continue                      # "Prof." on its own is not a person
        hit = near = None
        for i, (prev_core, _prev) in enumerate(out):
            if core == prev_core:
                hit = i
                break
            r = difflib.SequenceMatcher(None, core, prev_core).ratio()
            if r >= 0.90:
                hit = i
                break
            if r >= 0.75 and near is None:
                near = i
        if hit is not None:
            if len(n) > len(out[hit][1]):
                out[hit] = (out[hit][0], n)   # keep the fullest spelling
            continue
        if near is not None:
            warn.append(f"\"{out[near][1]}\" dan \"{n}\" mungkin orang yang sama - "
                        f"sila semak ejaan nama.")
        out.append((core, n))
    return [n for _, n in out], warn


def normalise_analysis(data):
    """Force a model's output into the exact shape the document generators
    expect, whatever it actually produced.

    Grammar-constrained JSON guarantees the shape, but not every model or
    provider supports it - Groq's llama-3.3-70b rejects `json_schema` and falls
    back to free-form JSON. A live meeting died here because `slides` came back
    as a LIST instead of an object and gen_pptx called .get() on it. Coercing
    once, here, protects every caller: local, cloud, and the edit-then-generate
    flow.
    """
    if not isinstance(data, dict):
        data = {}

    for f in ("meeting_title", "date", "time", "location"):
        data[f] = _as_text(data.get(f)).strip()

    data["attendees"], _name_warnings = _clean_attendees(
        [_as_text(a).strip() for a in _as_list(data.get("attendees")) if _as_text(a).strip()])
    for f in ("activities", "key_points", "key_takeaways", "important_notes"):
        data[f] = [_as_text(x).strip() for x in _as_list(data.get(f)) if _as_text(x).strip()]
    # Names that look like the same person go to the reader, not silently into
    # the attendance list.
    for _w in _name_warnings:
        if _w not in data["important_notes"]:
            data["important_notes"].append(_w)

    agenda = []
    for it in _as_list(data.get("agenda_items")):
        if isinstance(it, dict):
            agenda.append({"topic": _as_text(it.get("topic")).strip(),
                           "discussion": _as_text(it.get("discussion")).strip(),
                           "decision": _as_text(it.get("decision")).strip()})
        elif _as_text(it).strip():
            agenda.append({"topic": _as_text(it).strip(), "discussion": "", "decision": ""})
    data["agenda_items"] = agenda

    actions = []
    for it in _as_list(data.get("action_items")):
        if isinstance(it, dict):
            _task = _as_text(it.get("task")).strip()
            if not _task:
                continue            # an empty task is a blank row, not an action
            actions.append({"task": _task,
                            "owner": _as_text(it.get("owner")).strip(),
                            "deadline": _as_text(it.get("deadline")).strip()})
        elif _as_text(it).strip():
            actions.append({"task": _as_text(it).strip(), "owner": "", "deadline": ""})
    data["action_items"] = actions

    theme = _as_dict(data.get("theme"))
    data["theme"] = {
        "primary_hex": (_as_text(theme.get("primary_hex")).strip().lstrip("#") or "1E2761")[:6],
        "accent_hex": (_as_text(theme.get("accent_hex")).strip().lstrip("#") or "FFD500")[:6],
        "mood": _as_text(theme.get("mood")).strip(),
    }

    # slides: the field that actually broke a live meeting.
    slides = data.get("slides")
    if isinstance(slides, list):
        # A bare list of slides, or [{title_slide:..}, {content_slides:..}].
        merged = {}
        loose = []
        for entry in slides:
            if isinstance(entry, dict) and ("title_slide" in entry or "content_slides" in entry):
                merged.update(entry)
            else:
                loose.append(entry)
        if loose and "content_slides" not in merged:
            merged["content_slides"] = loose
        slides = merged
    slides = _as_dict(slides)

    title_slide = _as_dict(slides.get("title_slide"))
    content = []
    for sl in _as_list(slides.get("content_slides")):
        if isinstance(sl, dict):
            content.append({"heading": _as_text(sl.get("heading")).strip(),
                            "bullets": [_as_text(b).strip()
                                        for b in _as_list(sl.get("bullets"))
                                        if _as_text(b).strip()]})
        elif _as_text(sl).strip():
            content.append({"heading": _as_text(sl).strip(), "bullets": []})
    data["slides"] = {
        "title_slide": {"title": _as_text(title_slide.get("title")).strip()
                                 or data["meeting_title"] or "Meeting Minutes",
                        "subtitle": _as_text(title_slide.get("subtitle")).strip()
                                    or data.get("date", "")},
        "content_slides": content,
    }
    return data


def gen_docx(data, path):
    data = normalise_analysis(data)
    from docx.shared import Pt, RGBColor as _RGB
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    theme = data.get("theme", {})
    primary = hex_to_docx_rgb(theme.get("primary_hex", "1E2761"))
    ph = (theme.get("primary_hex", "1E2761") or "1E2761").lstrip("#")

    _tpl = os.path.join(DATA_DIR, "template.docx")
    doc = Document(_tpl) if os.path.exists(_tpl) else Document()
    _add_logo(doc)

    def shade(cell, hexcol):
        tcPr = cell._tc.get_or_add_tcPr()
        sh = OxmlElement("w:shd")
        sh.set(qn("w:val"), "clear"); sh.set(qn("w:fill"), hexcol)
        tcPr.append(sh)

    def section(label):
        p = doc.add_paragraph()
        p.space_before = Pt(14)
        run = p.add_run(label.upper())
        run.bold = True
        run.font.size = Pt(11)
        run.font.color.rgb = primary
        run.font.name = "Calibri"
        # bottom border on the section label
        pPr = p._p.get_or_add_pPr()
        pbdr = OxmlElement("w:pBdr")
        bottom = OxmlElement("w:bottom")
        bottom.set(qn("w:val"), "single"); bottom.set(qn("w:sz"), "6")
        bottom.set(qn("w:space"), "4"); bottom.set(qn("w:color"), ph)
        pbdr.append(bottom); pPr.append(pbdr)
        return p

    def bullet(text):
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(text)

    # ---- Header band (colored title) ----
    title_tbl = doc.add_table(rows=1, cols=1)
    tc = title_tbl.rows[0].cells[0]
    shade(tc, ph)
    tp = tc.paragraphs[0]
    tr = tp.add_run(data.get("meeting_title") or "Meeting Minutes")
    tr.bold = True; tr.font.size = Pt(20); tr.font.color.rgb = _RGB(0xFF, 0xFF, 0xFF)
    # meta line inside band
    mp = tc.add_paragraph()
    # Official minutes carry date, time and venue whether or not anyone said
    # them aloud. A blank to fill in beats a missing line.
    _lbl = doc_labels(data)
    meta = " · ".join(
        f"{_lbl[k]}: {data.get(k) or '________'}" for k in ("date", "time", "location"))
    mr = mp.add_run(meta)
    mr.font.size = Pt(10); mr.font.color.rgb = _RGB(0xE8, 0xEC, 0xF5)

    doc.add_paragraph()  # spacer

    L = doc_labels(data)

    # ---- Attendees ----
    if data.get("attendees"):
        section(L["attendees"])
        p = doc.add_paragraph()
        p.add_run(", ".join(data["attendees"]))

    # ---- Activities ----
    if data.get("activities"):
        section(L["activities"])
        for a in data["activities"]:
            bullet(a)

    # ---- Agenda / discussion ----
    if data.get("agenda_items"):
        section(L["matters"])
        for _n, item in enumerate(data["agenda_items"], 1):
            hp = doc.add_paragraph()
            # Numbered so a committee can refer to "perkara 4.0" in the next
            # meeting instead of quoting the heading back.
            hr = hp.add_run(f"{_n}.0  " + item.get("topic", ""))
            hr.bold = True; hr.font.size = Pt(12); hr.font.color.rgb = primary
            if item.get("discussion"):
                doc.add_paragraph(item["discussion"])
            if item.get("decision"):
                # The model is told to write the English phrase when nothing was
                # agreed; translate it so a Malay document stays Malay.
                _dec = item["decision"]
                if _dec.strip().rstrip(".").lower() == "no decision recorded":
                    _dec = L.get("no_decision", _dec)
                dp = doc.add_paragraph()
                dr = dp.add_run(L["decision"])
                dr.bold = True; dr.font.color.rgb = primary
                dp.add_run(_dec)

    # Key points and takeaways are deliberately NOT in the Word document.
    # The model fills every field it is given, so with one meeting's material
    # they came back as the agenda restated twice more in different grammar -
    # one real meeting produced 45 bullets carrying about 15 facts. They still
    # earn their place in the slides, where a condensed view is the point.

    # ---- Important notes ----
    if data.get("important_notes"):
        section(L["notes"])
        for n in data["important_notes"]:
            bullet(n)

    # ---- Action items table ----
    if data.get("action_items"):
        section(L["actions"])
        table = doc.add_table(rows=1, cols=3)
        try:
            table.style = "Table Grid"
        except KeyError:
            pass
        hdr = table.rows[0].cells
        for i, htext in enumerate([L["task"], L["owner"], L["deadline"]]):
            shade(hdr[i], ph)
            hp = hdr[i].paragraphs[0]
            hr = hp.add_run(htext)
            hr.bold = True; hr.font.color.rgb = _RGB(0xFF, 0xFF, 0xFF); hr.font.size = Pt(10)
        for idx, act in enumerate(data["action_items"]):
            row = table.add_row().cells
            row[0].text = act.get("task", "")
            row[1].text = act.get("owner", "")
            row[2].text = act.get("deadline", "")
            if idx % 2 == 0:
                for c in row:
                    shade(c, "F2F5FB")

    # ---- Signature block ----
    # Minutes are not official until someone puts their name to them.
    sig = doc.add_table(rows=1, cols=2)
    sig.autofit = True
    for cell, key in zip(sig.rows[0].cells, ("prepared_by", "verified_by")):
        cp = cell.paragraphs[0]
        cr = cp.add_run(L[key] + ":")
        cr.bold = True; cr.font.size = Pt(10); cr.font.color.rgb = primary
        cp.space_before = Pt(26)
        for line in ("", "........................................",
                     f"{L['name_line']}: ", f"{L['sig_date']}: "):
            lp = cell.add_paragraph()
            lr = lp.add_run(line)
            lr.font.size = Pt(10)

    # ---- Footer ----
    fp = doc.add_paragraph()
    fp.space_before = Pt(20)
    fr = fp.add_run("Generated by MinitAI · From voice to minutes, instantly")
    fr.font.size = Pt(8); fr.font.color.rgb = _RGB(0x8B, 0x94, 0xA9)
    fr.italic = True
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.save(path)


def gen_pptx(data, path):
    data = normalise_analysis(data)
    theme = data.get("theme", {})
    primary = hex_to_rgb(theme.get("primary_hex", "1E2761"), NAVY)
    accent = hex_to_rgb(theme.get("accent_hex", "FFD500"), WHITE)

    _tpl = os.path.join(DATA_DIR, "template.pptx")
    prs = Presentation(_tpl) if os.path.exists(_tpl) else Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(5.63)

    ts = data.get("slides", {}).get("title_slide", {})
    subtitle = ts.get("subtitle") or " | ".join(filter(None, [data.get("date"), data.get("time"), data.get("location")]))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = primary
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.55), Inches(10), Pt(4))
    accent_bar.fill.solid(); accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.7), Inches(2.7), Inches(8.6), Inches(1.5))
    p = box.text_frame.paragraphs[0]
    p.text = ts.get("title", "Meeting Summary")
    p.font.size, p.font.bold, p.font.color.rgb = Pt(40), True, WHITE
    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(8.6), Inches(1.0))
        sp = sbox.text_frame.paragraphs[0]
        sp.text, sp.font.size, sp.font.color.rgb = subtitle, Pt(18), WHITE

    for sd in data.get("slides", {}).get("content_slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = WHITE
        tbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.9))
        tp = tbox.text_frame.paragraphs[0]
        tp.text, tp.font.size, tp.font.bold, tp.font.color.rgb = sd.get("heading",""), Pt(30), True, primary
        bbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(8.4), Inches(5.0))
        tf = bbox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(sd.get("bullets", [])):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = f"•  {bullet}"
            para.font.size, para.font.color.rgb, para.space_after = Pt(18), DARK, Pt(12)

    action_items = data.get("action_items", [])
    if action_items:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = WHITE
        tbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.9))
        tp = tbox.text_frame.paragraphs[0]
        tp.text, tp.font.size, tp.font.bold, tp.font.color.rgb = "Tindakan", Pt(30), True, primary
        bbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(8.4), Inches(5.0))
        tf = bbox.text_frame
        tf.word_wrap = True
        for i, act in enumerate(action_items):
            line = f"{act.get('task','')} — {act.get('owner','')} ({act.get('deadline','')})"
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = f"•  {line}"
            para.font.size, para.font.color.rgb, para.space_after = Pt(18), DARK, Pt(12)

    prs.save(path)


import re
from datetime import datetime

def safe_filename(text, fallback):
    text = (text or "").strip()
    if not text:
        text = fallback
    text = re.sub(r'[<>:"/\\|?*]', '', text)
    text = re.sub(r'\s+', '_', text)
    text = text[:60].strip('_')
    return text or "meeting"


def unique_path(path):
    if not os.path.exists(path):
        return path
    base, ext = os.path.splitext(path)
    i = 2
    while os.path.exists(f"{base}_{i}{ext}"):
        i += 1
    return f"{base}_{i}{ext}"



PDF_UNAVAILABLE_REASON = ""       # read by the UI so the user is told why


def gen_pdf(docx_path):
    """Convert docx -> pdf. This needs Microsoft Word installed, because
    docx2pdf drives Word itself. Plenty of machines have no Word (WPS,
    LibreOffice, Google Docs), and this used to no-op in silence: the user got
    no PDF and no explanation. Now the reason is recorded and shown."""
    global PDF_UNAVAILABLE_REASON
    PDF_UNAVAILABLE_REASON = ""
    if not _cfg.get("export_pdf", True):
        return None
    try:
        from docx2pdf import convert
    except ImportError:
        PDF_UNAVAILABLE_REASON = ("PDF export needs the docx2pdf package. "
                                  "Re-run SETUP to install it.")
        logging.info("PDF export skipped: docx2pdf not installed")
        return None
    try:
        pdf_path = docx_path.rsplit(".", 1)[0] + ".pdf"
        convert(docx_path, pdf_path)
        if os.path.exists(pdf_path):
            return pdf_path
        PDF_UNAVAILABLE_REASON = ("PDF export produced nothing. Microsoft Word "
                                  "must be installed for PDF export.")
        return None
    except Exception as e:
        PDF_UNAVAILABLE_REASON = ("PDF export needs Microsoft Word on this PC. "
                                  "Your Word file was still created normally.")
        logging.info(f"PDF export skipped: {type(e).__name__}")
        return None

def write_document(doc_type, instructions, lang=None, tone="professional"):
    """Generate structured content for reports, letters, emails, notices via Ollama.
    Returns a dict shaped per doc_type."""
    ensure_ollama_ready()

    lang_rule = _LANG_RULES.get(lang or "auto", _LANG_RULES["auto"])

    schemas = {
        "report": '{"title":"","date":"","author":"","sections":[{"heading":"","body":""}],"summary":""}',
        "letter": '{"sender":"","sender_address":"","recipient":"","recipient_address":"","date":"","subject":"","salutation":"","body_paragraphs":[""],"closing":"","signature":""}',
        "email": '{"to":"","subject":"","greeting":"","body_paragraphs":[""],"signoff":"","sender":""}',
        "notice": '{"organization":"","title":"","reference_no":"","date":"","body_paragraphs":[""],"issued_by":"","position":""}',
    }
    schema = schemas.get(doc_type, schemas["report"])

    guidance = {
        "report": "Write a formal report. Organize into clear sections with headings. Include a brief summary.",
        "letter": "Write a formal official letter with proper structure (sender, recipient, date, subject, salutation, body, closing, signature).",
        "email": "Write a professional email. Keep it clear and appropriately concise.",
        "notice": "Write an official notice/circular suitable for posting. Formal and clear.",
    }.get(doc_type, "Write a formal document.")

    system = f"""{lang_rule} You are an expert administrative writer for a campus/organization.
{guidance}
Tone: {tone}. Fill realistic details from the user's instructions; leave truly unknown fields as empty strings.
Output ONLY valid JSON, no markdown, matching exactly this shape:
{schema}"""

    _payload = {"model": OLLAMA_MODEL, "system": system, "prompt": instructions,
                "stream": False, "format": "json"}
    if _cfg.get("low_resource_mode", True):
        _payload["keep_alive"] = _KEEP_ALIVE
    resp = requests.post(OLLAMA_URL + "/api/generate", timeout=300, json=_payload)
    raw = resp.json()["response"]
    return _extract_json(raw, schema)


def _extract_json(raw, schema):
    """Parse JSON from possibly-messy AI output."""
    import re
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        pass
    # try to find the first {...} block
    m = re.search(r'\{.*\}', raw, re.DOTALL)
    if m:
        try:
            return json.loads(m.group(0))
        except json.JSONDecodeError:
            pass
    # last resort: return the schema shape with the raw text as body so user gets something
    try:
        shape = json.loads(schema)
    except Exception:
        shape = {}
    if "sections" in shape:
        shape["title"] = "Document"
        shape["sections"] = [{"heading": "Content", "body": raw.strip()[:4000]}]
    elif "body_paragraphs" in shape:
        shape["body_paragraphs"] = [raw.strip()[:4000]]
    return shape


def _add_logo(doc):
    """Insert brand logo at top of document if user uploaded one."""
    import glob
    from docx.shared import Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    for e in (".png", ".jpg", ".jpeg"):
        lp = os.path.join(DATA_DIR, "brand_logo" + e)
        if os.path.exists(lp):
            try:
                p = doc.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                p.add_run().add_picture(lp, width=Inches(1.8))
                return True
            except Exception:
                return False
    return False


def _doc_base(path, template_name=None):
    from docx import Document as _D
    tpl = os.path.join(DATA_DIR, template_name) if template_name else None
    if tpl and os.path.exists(tpl):
        return _D(tpl)
    return _D()


def gen_report_docx(data, path):
    from docx.shared import Pt, RGBColor as _RGB
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    doc = _doc_base(path, "template_report.docx")
    _add_logo(doc)
    navy = _RGB(0x1E, 0x3A, 0x5F)

    t = doc.add_paragraph()
    tr = t.add_run(data.get("title") or "Report")
    tr.bold = True; tr.font.size = Pt(20); tr.font.color.rgb = navy
    meta = doc.add_paragraph()
    mparts = [x for x in [data.get("date"), ("By " + data["author"]) if data.get("author") else ""] if x]
    mr = meta.add_run("  ·  ".join(mparts))
    mr.font.size = Pt(10); mr.font.color.rgb = _RGB(0x6B, 0x72, 0x80)

    if data.get("summary"):
        h = doc.add_paragraph(); hr = h.add_run("EXECUTIVE SUMMARY")
        hr.bold = True; hr.font.size = Pt(11); hr.font.color.rgb = navy
        doc.add_paragraph(data["summary"])

    for sec in data.get("sections", []):
        h = doc.add_paragraph(); hr = h.add_run((sec.get("heading") or "").upper())
        hr.bold = True; hr.font.size = Pt(12); hr.font.color.rgb = navy
        if sec.get("body"):
            doc.add_paragraph(sec["body"])

    _footer(doc)
    doc.save(path)


def gen_letter_docx(data, path):
    from docx.shared import Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    doc = _doc_base(path, "template_letter.docx")
    _add_logo(doc)

    if data.get("sender"):
        doc.add_paragraph(data["sender"])
    if data.get("sender_address"):
        for line in str(data["sender_address"]).split("\n"):
            doc.add_paragraph(line)
    if data.get("date"):
        p = doc.add_paragraph(data["date"]); p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    doc.add_paragraph()
    if data.get("recipient"):
        doc.add_paragraph(data["recipient"])
    if data.get("recipient_address"):
        for line in str(data["recipient_address"]).split("\n"):
            doc.add_paragraph(line)
    doc.add_paragraph()
    if data.get("subject"):
        sp = doc.add_paragraph(); sr = sp.add_run(data["subject"]); sr.bold = True
    doc.add_paragraph()
    if data.get("salutation"):
        doc.add_paragraph(data["salutation"])
    for para in data.get("body_paragraphs", []):
        doc.add_paragraph(para)
    doc.add_paragraph()
    if data.get("closing"):
        doc.add_paragraph(data["closing"])
    if data.get("signature"):
        doc.add_paragraph(); doc.add_paragraph(data["signature"])

    doc.save(path)


def gen_notice_docx(data, path):
    from docx.shared import Pt, RGBColor as _RGB
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    doc = _doc_base(path, "template_notice.docx")
    _add_logo(doc)
    navy = _RGB(0x1E, 0x3A, 0x5F)

    if data.get("organization"):
        o = doc.add_paragraph(); o.alignment = WD_ALIGN_PARAGRAPH.CENTER
        orn = o.add_run(data["organization"]); orn.bold = True; orn.font.size = Pt(14); orn.font.color.rgb = navy

    tt = doc.add_paragraph(); tt.alignment = WD_ALIGN_PARAGRAPH.CENTER
    ttr = tt.add_run((data.get("title") or "NOTICE").upper()); ttr.bold = True; ttr.font.size = Pt(16)

    meta = doc.add_paragraph()
    mparts = [x for x in [("Ref: " + data["reference_no"]) if data.get("reference_no") else "", data.get("date")] if x]
    if mparts:
        meta.add_run("     ".join(mparts)).font.size = Pt(10)

    doc.add_paragraph()
    for para in data.get("body_paragraphs", []):
        doc.add_paragraph(para)

    doc.add_paragraph()
    if data.get("issued_by"):
        doc.add_paragraph(data["issued_by"])
    if data.get("position"):
        doc.add_paragraph(data["position"])

    doc.save(path)


def _footer(doc):
    from docx.shared import Pt, RGBColor as _RGB
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    fp = doc.add_paragraph(); fp.space_before = Pt(18)
    fr = fp.add_run("Generated by MinitAI · From voice to minutes, instantly")
    fr.font.size = Pt(8); fr.italic = True; fr.font.color.rgb = _RGB(0x8B, 0x94, 0xA9)
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER


def gen_write_output(doc_type, data, path):
    """Route to the right docx generator."""
    if doc_type == "report":
        gen_report_docx(data, path)
    elif doc_type == "letter":
        gen_letter_docx(data, path)
    elif doc_type == "notice":
        gen_notice_docx(data, path)
    elif doc_type == "email":
        # email as a simple docx (also returned as plain text in the route)
        from docx.shared import Pt
        doc = _doc_base(path)
        if data.get("subject"):
            p = doc.add_paragraph(); r = p.add_run("Subject: " + data["subject"]); r.bold = True
        doc.add_paragraph(data.get("greeting", ""))
        for para in data.get("body_paragraphs", []):
            doc.add_paragraph(para)
        doc.add_paragraph(data.get("signoff", ""))
        doc.add_paragraph(data.get("sender", ""))
        doc.save(path)
    else:
        gen_report_docx(data, path)


def pdf_from_images(image_paths, out_path):
    """Combine images (JPG/PNG) into a single PDF, one image per page."""
    from PIL import Image
    imgs = []
    for p in image_paths:
        im = Image.open(p)
        if im.mode in ("RGBA", "P", "LA"):
            im = im.convert("RGB")
        imgs.append(im)
    if not imgs:
        raise RuntimeError("No valid images provided.")
    imgs[0].save(out_path, "PDF", save_all=True, append_images=imgs[1:])
    return out_path


def pdf_merge(pdf_paths, out_path):
    """Merge multiple PDFs into one."""
    from pypdf import PdfWriter, PdfReader
    writer = PdfWriter()
    for p in pdf_paths:
        reader = PdfReader(p)
        for page in reader.pages:
            writer.add_page(page)
    with open(out_path, "wb") as f:
        writer.write(f)
    return out_path


def pdf_split(pdf_path, out_dir, ranges=None):
    """Split a PDF. If ranges given (list of (start,end) 1-indexed inclusive),
    produce one file per range; otherwise one file per page. Returns list of paths."""
    from pypdf import PdfReader, PdfWriter
    reader = PdfReader(pdf_path)
    n = len(reader.pages)
    base = os.path.splitext(os.path.basename(pdf_path))[0]
    out = []
    if ranges:
        for (a, b) in ranges:
            a = max(1, a); b = min(n, b)
            if a > b:
                continue
            writer = PdfWriter()
            for i in range(a - 1, b):
                writer.add_page(reader.pages[i])
            op = unique_path(os.path.join(out_dir, f"{base}_pages_{a}-{b}.pdf"))
            with open(op, "wb") as f:
                writer.write(f)
            out.append(op)
    else:
        for i in range(n):
            writer = PdfWriter()
            writer.add_page(reader.pages[i])
            op = unique_path(os.path.join(out_dir, f"{base}_page_{i+1}.pdf"))
            with open(op, "wb") as f:
                writer.write(f)
            out.append(op)
    return out


def word_to_pdf(docx_path, out_path=None):
    """Convert a .docx to PDF."""
    if out_path is None:
        out_path = os.path.splitext(docx_path)[0] + ".pdf"
    try:
        from docx2pdf import convert
        convert(docx_path, out_path)
        if os.path.exists(out_path):
            return out_path
    except Exception:
        pass
    import subprocess, shutil
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice:
        outdir = os.path.dirname(out_path) or "."
        subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir", outdir, docx_path],
                       timeout=120, capture_output=True)
        produced = os.path.join(outdir, os.path.splitext(os.path.basename(docx_path))[0] + ".pdf")
        if os.path.exists(produced):
            if produced != out_path:
                shutil.move(produced, out_path)
            return out_path
    raise RuntimeError("PDF conversion needs Microsoft Word or LibreOffice installed.")


def pdf_page_count(pdf_path):
    from pypdf import PdfReader
    return len(PdfReader(pdf_path).pages)


def extract_text_from_file(path):
    """Extract plain text from an uploaded document: PDF, PowerPoint, Word, or txt.
    Used by Study tools so students can generate from their slides."""
    ext = os.path.splitext(path)[1].lower()
    text = ""
    try:
        if ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(path)
            text = "\n".join((pg.extract_text() or "") for pg in reader.pages)
        elif ext in (".pptx", ".ppt"):
            from pptx import Presentation
            prs = Presentation(path)
            chunks = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = "".join(run.text for run in para.runs)
                            if line.strip():
                                chunks.append(line)
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [c.text for c in row.cells]
                            chunks.append(" | ".join(cells))
            text = "\n".join(chunks)
        elif ext == ".docx":
            from docx import Document as _D
            doc = _D(path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext in (".txt", ".md"):
            with open(path, encoding="utf-8", errors="ignore") as f:
                text = f.read()
        else:
            raise RuntimeError(f"Unsupported file type: {ext}. Use PDF, PowerPoint, Word, or text.")
    except Exception as e:
        raise RuntimeError(f"Couldn't read that file: {e}")
    text = text.strip()
    if len(text) < 20:
        raise RuntimeError("Not enough readable text found. If your slides are scanned images, text can't be extracted.")
    return text[:20000]


def _study_source_text(source_text=None):
    """Get text to study from: explicit text, or latest transcript."""
    if source_text and source_text.strip():
        return source_text.strip()[:20000]
    import glob
    txts = sorted(glob.glob(os.path.join(OUTPUT_DIR, "*transcript*.txt")), key=os.path.getmtime, reverse=True)
    if txts:
        with open(txts[0], encoding="utf-8") as f:
            return f.read()[:20000]
    return ""


def gen_flashcards(source_text=None, lang=None, count=10):
    """Generate study flashcards (question/answer pairs) from lecture content."""
    ensure_ollama_ready()
    text = _study_source_text(source_text)
    if not text:
        raise RuntimeError("No lecture content found. Record or upload a lecture first.")
    lang_rule = _LANG_RULES.get(lang or "auto", _LANG_RULES.get("auto", ""))
    schema = '{"cards":[{"front":"","back":""}]}'
    system = f"""{lang_rule} You are a study assistant creating flashcards for a university student.
From the lecture content, create {count} flashcards covering the most important concepts, definitions, and facts.
Front = a clear question or term. Back = a concise, accurate answer.
Output ONLY valid JSON matching: {schema}"""
    payload = {"model": OLLAMA_MODEL, "system": system, "prompt": text, "stream": False, "format": "json"}
    if _cfg.get("low_resource_mode", True):
        payload["keep_alive"] = _KEEP_ALIVE
    resp = requests.post(OLLAMA_URL + "/api/generate", timeout=300, json=payload)
    data = _extract_json(resp.json()["response"], schema)
    return data.get("cards", [])


def gen_quiz(source_text=None, lang=None, count=8, qtypes=None):
    """Generate a quiz from lecture content. qtypes = list of 'mcq','tf','short'."""
    ensure_ollama_ready()
    text = _study_source_text(source_text)
    if not text:
        raise RuntimeError("No lecture content found. Record or upload a lecture first.")
    lang_rule = _LANG_RULES.get(lang or "auto", _LANG_RULES.get("auto", ""))
    types_txt = ", ".join(qtypes) if qtypes else "mcq, tf, short"
    schema = ('{"questions":[{"type":"mcq","question":"","options":["","","",""],'
              '"answer":"","explanation":""}]}')
    system = f"""{lang_rule} You are a university exam-prep assistant.
Create a {count}-question quiz from the lecture content. Use these question types: {types_txt}.
For "mcq": include 4 options and set answer to the correct option text.
For "tf": options must be ["True","False"] and answer one of them.
For "short": omit options, answer is a brief model answer.
Always include a one-sentence explanation.
Output ONLY valid JSON matching: {schema}"""
    payload = {"model": OLLAMA_MODEL, "system": system, "prompt": text, "stream": False, "format": "json"}
    if _cfg.get("low_resource_mode", True):
        payload["keep_alive"] = _KEEP_ALIVE
    resp = requests.post(OLLAMA_URL + "/api/generate", timeout=300, json=payload)
    data = _extract_json(resp.json()["response"], schema)
    return data.get("questions", [])


def gen_study_notes(source_text=None, lang=None, depth="medium"):
    """Generate study notes at a chosen depth: quick, medium, detailed."""
    ensure_ollama_ready()
    text = _study_source_text(source_text)
    if not text:
        raise RuntimeError("No lecture content found. Record or upload a lecture first.")
    lang_rule = _LANG_RULES.get(lang or "auto", _LANG_RULES.get("auto", ""))
    depth_rule = {
        "quick": "Write a QUICK summary: the 5-7 most important points only, as short bullet points.",
        "medium": "Write MEDIUM study notes: organized by topic with headings, key points, and important definitions.",
        "detailed": "Write DETAILED study notes: full coverage with headings, explanations, definitions, examples, and any formulas.",
    }.get(depth, "Write clear study notes.")
    schema = '{"title":"","sections":[{"heading":"","points":[""]}],"key_terms":[{"term":"","definition":""}]}'
    system = f"""{lang_rule} You are a study-notes assistant for a university student.
{depth_rule}
Output ONLY valid JSON matching: {schema}"""
    payload = {"model": OLLAMA_MODEL, "system": system, "prompt": text, "stream": False, "format": "json"}
    if _cfg.get("low_resource_mode", True):
        payload["keep_alive"] = _KEEP_ALIVE
    resp = requests.post(OLLAMA_URL + "/api/generate", timeout=300, json=payload)
    return _extract_json(resp.json()["response"], schema)


def gen_flashcards_docx(cards, path):
    """Export flashcards as a printable Word doc (foldable cards)."""
    from docx.shared import Pt, RGBColor as _RGB
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    doc = _doc_base(path)
    _add_logo(doc)
    navy = _RGB(0x1E, 0x3A, 0x5F)
    t = doc.add_paragraph(); tr = t.add_run("Flashcards")
    tr.bold = True; tr.font.size = Pt(20); tr.font.color.rgb = navy
    for i, c in enumerate(cards, 1):
        q = doc.add_paragraph()
        qr = q.add_run(f"Q{i}. {c.get('front','')}"); qr.bold = True; qr.font.size = Pt(12)
        a = doc.add_paragraph()
        ar = a.add_run(f"A: {c.get('back','')}"); ar.font.size = Pt(11); ar.font.color.rgb = _RGB(0x33,0x38,0x4a)
        doc.add_paragraph()
    doc.save(path)
    return path


def gen_quiz_docx(questions, path, with_answers=True):
    """Export quiz as a Word doc. Questions first, answer key at the end."""
    from docx.shared import Pt, RGBColor as _RGB
    doc = _doc_base(path)
    _add_logo(doc)
    navy = _RGB(0x1E, 0x3A, 0x5F)
    t = doc.add_paragraph(); tr = t.add_run("Quiz")
    tr.bold = True; tr.font.size = Pt(20); tr.font.color.rgb = navy
    for i, q in enumerate(questions, 1):
        p = doc.add_paragraph()
        pr = p.add_run(f"{i}. {q.get('question','')}"); pr.bold = True; pr.font.size = Pt(12)
        opts = q.get("options") or []
        for j, opt in enumerate(opts):
            letter = chr(65 + j)
            doc.add_paragraph(f"   {letter}. {opt}")
        doc.add_paragraph()
    if with_answers:
        doc.add_page_break()
        ap = doc.add_paragraph(); ar = ap.add_run("Answer Key")
        ar.bold = True; ar.font.size = Pt(16); ar.font.color.rgb = navy
        for i, q in enumerate(questions, 1):
            a = doc.add_paragraph()
            a.add_run(f"{i}. ").bold = True
            a.add_run(str(q.get("answer", "")))
            if q.get("explanation"):
                e = doc.add_paragraph(f"   → {q['explanation']}")
                for run in e.runs:
                    run.font.size = Pt(10); run.font.color.rgb = _RGB(0x6b,0x72,0x80)
    doc.save(path)
    return path


def build_diagnostics(max_log_lines=400):
    """Build a single shareable report describing this installation and any errors.

    PRIVACY: deliberately excludes meeting transcripts, document contents and
    file names of user documents — only technical facts needed for debugging.
    """
    import platform, shutil as _sh
    L = []
    def add(s=""):
        L.append(str(s))

    add("=" * 62)
    add("  MinitAI DIAGNOSTIC REPORT")
    add("  Safe to share: contains NO meeting content or transcripts.")
    add("=" * 62)
    add(f"Generated : {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    add(f"Version   : {APP_VERSION}")
    add()

    # ---------- system ----------
    add("-" * 62)
    add("SYSTEM")
    add("-" * 62)
    try:
        add(f"OS            : {platform.system()} {platform.release()} ({platform.version()})")
        add(f"Machine       : {platform.machine()}")
        add(f"Python        : {sys.version.split()[0]}  ({sys.executable})")
    except Exception as e:
        add(f"OS info failed: {e}")
    try:
        import psutil
        vm = psutil.virtual_memory()
        add(f"RAM total     : {vm.total/1024**3:.1f} GB")
        add(f"RAM available : {vm.available/1024**3:.1f} GB")
        add(f"CPU cores     : {psutil.cpu_count(logical=True)}")
    except Exception as e:
        add(f"RAM/CPU info failed: {e}")
    try:
        du = _sh.disk_usage(DATA_DIR)
        add(f"Disk free     : {du.free/1024**3:.1f} GB of {du.total/1024**3:.1f} GB")
    except Exception as e:
        add(f"Disk info failed: {e}")
    add()

    # ---------- chosen settings ----------
    add("-" * 62)
    add("MINITAI SETTINGS (auto-detected)")
    add("-" * 62)
    add(f"Hardware tier : {_tier}")
    add(f"Whisper model : {MODEL_SIZE}")
    add(f"Beam size     : {BEAM_SIZE}")
    add(f"Spoken lang   : {TRANSCRIBE_LANG or 'detect once, then lock'}")
    add(f"Name hints    : {'set' if INITIAL_PROMPT else 'none'}")
    # Engine status. NEVER prints the API key itself - only whether one exists.
    try:
        add(f"Engine mode   : {ENGINE_MODE}")
        add(f"Always private: {FORCE_LOCAL_ALWAYS}")
        add(f"Measured speed: {_load_rtf():.2f}x realtime "
            f"(1 hour of audio ~ {_load_rtf() * 60:.0f} min on this PC)")
        if _cloud:
            _k = _cloud.get_key(DATA_DIR)
            add(f"Online key    : {'present (' + str(len(_k)) + ' chars)' if _k else 'not set up'}")
            add(f"Online reachable: {_cloud.available(DATA_DIR) if _k else False}")
        else:
            add("Online engine : not installed")
        _e, _why = choose_engine(3600)
        add(f"A 1-hour meeting would run: {_e} ({_why})")
    except Exception as _e:
        add(f"Engine status : unavailable ({type(_e).__name__})")
    try:
        add(f"PDF export    : {'Word not available - ' + PDF_UNAVAILABLE_REASON if PDF_UNAVAILABLE_REASON else 'ok / untested'}")
    except Exception:
        pass
    add(f"AI model      : {OLLAMA_MODEL}")
    add(f"CPU threads   : {_threads}")
    add(f"Model keep    : {_KEEP_ALIVE}")
    add(f"Free whisper  : {_FREE_WHISPER_BETWEEN}")
    add(f"Data folder   : {DATA_DIR}")
    try:
        add("config.json   : " + json.dumps(_cfg, ensure_ascii=False))
    except Exception as e:
        add(f"config read failed: {e}")
    add()

    # ---------- python packages ----------
    add("-" * 62)
    add("PACKAGES")
    add("-" * 62)
    for mod in ("faster_whisper", "flask", "docx", "pptx", "pypdf", "PIL",
                "psutil", "requests", "numpy", "sounddevice", "imageio_ffmpeg"):
        try:
            m = __import__(mod)
            add(f"  {mod:18s} {getattr(m, '__version__', 'installed')}")
        except Exception as e:
            add(f"  {mod:18s} MISSING ({type(e).__name__})")
    add()

    # ---------- ollama ----------
    add("-" * 62)
    add("OLLAMA")
    add("-" * 62)
    try:
        r = requests.get(OLLAMA_URL + "/api/tags", timeout=5)
        models = [m.get("name", "?") for m in r.json().get("models", [])]
        add(f"  Status        : running ({OLLAMA_URL})")
        add(f"  Models present: {', '.join(models) if models else '(none)'}")
        add(f"  Required model: {OLLAMA_MODEL} "
            f"{'FOUND' if any(OLLAMA_MODEL in m for m in models) else '*** MISSING ***'}")
    except Exception as e:
        add(f"  Status        : NOT REACHABLE ({type(e).__name__}: {e})")
    add()

    # ---------- activity (counts only, no names/content) ----------
    add("-" * 62)
    add("ACTIVITY")
    add("-" * 62)
    try:
        n_out = len(os.listdir(OUTPUT_DIR)) if os.path.isdir(OUTPUT_DIR) else 0
        n_aud = len(os.listdir(AUDIO_DIR)) if os.path.isdir(AUDIO_DIR) else 0
        add(f"  Files in output folder : {n_out}")
        add(f"  Files in audio folder  : {n_aud}")
    except Exception as e:
        add(f"  folder scan failed: {e}")
    try:
        con = sqlite3.connect(DB_PATH)
        cnt = con.execute("SELECT COUNT(*) FROM meetings").fetchone()[0]
        con.close()
        add(f"  Meetings recorded      : {cnt}")
    except Exception as e:
        add(f"  database: not available ({type(e).__name__})")
    add()

    # ---------- errors ----------
    add("-" * 62)
    add("ERRORS AND WARNINGS (most recent last)")
    add("-" * 62)
    try:
        if os.path.exists(LOG_PATH):
            with open(LOG_PATH, encoding="utf-8", errors="ignore") as f:
                lines = f.readlines()
            problems = [l.rstrip() for l in lines
                        if "[ERROR]" in l or "[WARNING]" in l or "Traceback" in l]
            add(f"  ({len(problems)} problem line(s) found in the log)")
            add()
            for l in problems[-120:]:
                add("  " + l)
        else:
            add("  No log file yet.")
    except Exception as e:
        add(f"  could not read log: {e}")
    add()

    # ---------- recent log tail ----------
    add("-" * 62)
    add(f"FULL LOG (last {max_log_lines} lines)")
    add("-" * 62)
    try:
        if os.path.exists(LOG_PATH):
            with open(LOG_PATH, encoding="utf-8", errors="ignore") as f:
                for l in f.readlines()[-max_log_lines:]:
                    add("  " + l.rstrip())
        else:
            add("  (none)")
    except Exception as e:
        add(f"  could not read log: {e}")
    add()
    add("=" * 62)
    add("  END OF REPORT")
    add("=" * 62)
    return "\n".join(L)


def save_diagnostics():
    """Write the diagnostic report to the Desktop (and the data folder). Returns path."""
    name = f"MinitAI_Error_Report_{datetime.now().strftime('%Y-%m-%d_%H-%M')}.txt"
    text = build_diagnostics()
    primary = unique_path(os.path.join(OUTPUT_DIR, name))
    with open(primary, "w", encoding="utf-8") as f:
        f.write(text)
    try:
        save_to_desktop(primary)
    except Exception:
        pass
    return primary


def log_environment():
    """Record machine + settings once at startup so every log has context."""
    try:
        import platform, psutil
        vm = psutil.virtual_memory()
        logging.info(
            f"MinitAI {APP_VERSION} | {platform.system()} {platform.release()} | "
            f"RAM {vm.total/1024**3:.1f}GB (free {vm.available/1024**3:.1f}GB) | "
            f"cores {psutil.cpu_count(logical=True)} | tier {_tier} | "
            f"whisper {MODEL_SIZE} | llm {OLLAMA_MODEL} | threads {_threads}"
        )
    except Exception as e:
        logging.warning(f"environment logging failed: {e}")


def find_unprocessed_audio():
    """Return audio files that were saved but never produced a transcript
    (i.e. the app likely crashed mid-processing). Used for crash recovery."""
    import glob
    recoverable = []
    try:
        audio_files = []
        for ext in AUDIO_EXTS:
            audio_files += glob.glob(os.path.join(AUDIO_DIR, f"*{ext}"))
        for af in audio_files:
            base = os.path.splitext(os.path.basename(af))[0]
            # if any transcript starts with this base, it was processed
            matches = glob.glob(os.path.join(OUTPUT_DIR, f"{base}_*transcript*.txt"))
            if not matches and os.path.getsize(af) > 2048:
                recoverable.append({
                    "path": af,
                    "name": os.path.basename(af),
                    "size_mb": round(os.path.getsize(af) / (1024*1024), 1),
                    "when": datetime.fromtimestamp(os.path.getmtime(af)).strftime("%Y-%m-%d %H:%M"),
                })
    except Exception:
        pass
    return recoverable


def process(audio_path):
    base = os.path.splitext(os.path.basename(audio_path))[0]
    print(f"Processing {base}...")
    transcript_text = transcribe(audio_path)

    ts_early = datetime.now().strftime("%Y-%m-%d_%H-%M")
    with open(os.path.join(OUTPUT_DIR, f"{base}_{ts_early}_transcript.txt"), "w", encoding="utf-8") as f:
        f.write(transcript_text)

    data = analyze(transcript_text)

    ts = datetime.now().strftime("%Y-%m-%d_%H-%M")
    name = ts

    docx_path = unique_path(os.path.join(OUTPUT_DIR, f"{name}_minutes.docx"))
    pptx_path = unique_path(os.path.join(OUTPUT_DIR, f"{name}_slides.pptx"))
    gen_docx(data, docx_path)
    gen_pptx(data, pptx_path)
    gen_pdf(docx_path)
    log_meeting(data.get("meeting_title") or base, data.get("date",""), docx_path, pptx_path)
    print(f"Done: {os.path.basename(docx_path)}, {os.path.basename(pptx_path)}")


if __name__ == "__main__":
    # Imported here, not at module load: the folder watcher is a desktop-only
    # convenience and the hosted version does not install watchdog.
    from watchdog.observers import Observer
    from watchdog.events import FileSystemEventHandler

    class Handler(FileSystemEventHandler):
        def on_created(self, event):
            if event.is_directory:
                return
            if event.src_path.lower().endswith(AUDIO_EXTS):
                time.sleep(2)  # wait for file write to finish
                try:
                    process(event.src_path)
                except Exception as e:
                    print(f"Error processing {event.src_path}: {e}")

    os.makedirs(AUDIO_DIR, exist_ok=True)
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    observer = Observer()
    observer.schedule(Handler(), AUDIO_DIR, recursive=False)
    observer.start()
    print(f"Watching {AUDIO_DIR}/ — drop audio files to auto-process. Ctrl+C to stop.")
    try:
        while True:
            time.sleep(1)
    except KeyboardInterrupt:
        observer.stop()
    observer.join()
